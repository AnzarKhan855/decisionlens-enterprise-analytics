from __future__ import annotations

from typing import Any, Dict, List, Optional
from datetime import datetime
from pathlib import Path

from app.logging.logger import get_logger

logger = get_logger(__name__)


class PDFExporter:
    @staticmethod
    def export(report: Dict[str, Any], title: str = "DecisionLens Executive Report") -> bytes:
        try:
            from reportlab.lib import colors
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch

            buffer = []
            doc = SimpleDocTemplate(
                "/tmp/decisionlens_report.pdf",
                pagesize=letter,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=18,
            )

            styles = getSampleStyleSheet()
            title_style = ParagraphStyle(
                "CustomTitle",
                parent=styles["Heading1"],
                fontSize=24,
                textColor=colors.HexColor("#1a1a2e"),
                spaceAfter=30,
            )
            heading_style = ParagraphStyle(
                "CustomHeading",
                parent=styles["Heading2"],
                fontSize=16,
                textColor=colors.HexColor("#16213e"),
                spaceAfter=12,
                spaceBefore=12,
            )
            body_style = ParagraphStyle(
                "CustomBody",
                parent=styles["BodyText"],
                fontSize=10,
                leading=14,
            )

            story = []
            story.append(Paragraph(title, title_style))
            story.append(Spacer(1, 0.2 * inch))

            meta = [
                ["Generated", report.get("generated_at", "")[:19]],
                ["Domain", report.get("domain", "N/A")],
                ["Dataset Type", report.get("dataset_type", "N/A")],
            ]
            meta_table = Table(meta, colWidths=[1.5 * inch, 4 * inch])
            meta_table.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (0, -1), colors.HexColor("#e94560")),
                ("TEXTCOLOR", (0, 0), (0, -1), colors.whitesmoke),
                ("BACKGROUND", (1, 0), (1, -1), colors.HexColor("#f5f5f5")),
                ("TEXTCOLOR", (1, 0), (1, -1), colors.black),
                ("ALIGN", (0, 0), (-1, -1), "LEFT"),
                ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), 10),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("LEFTPADDING", (0, 0), (-1, -1), 8),
                ("RIGHTPADDING", (0, 0), (-1, -1), 8),
                ("TOPPADDING", (0, 0), (-1, -1), 6),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
            ]))
            story.append(meta_table)
            story.append(Spacer(1, 0.3 * inch))

            sections = report.get("sections", {})
            section_order = [
                "executive_summary", "business_health", "kpi_summary", "kpi_overview",
                "revenue", "customers", "products", "categories", "trend_analysis",
                "forecast", "root_cause_analysis", "risks", "opportunities",
                "recommendations", "evidence", "charts", "what_if_analysis",
                "recommended_actions", "business_impact", "key_findings",
                "confidence_evidence", "roadmap_30_90_180", "domain_specific",
            ]

            for section_name in section_order:
                section_data = sections.get(section_name)
                if section_data is None:
                    continue
                story.append(Paragraph(section_name.replace("_", " ").title(), heading_style))
                story.append(PDFExporter._render_section_pdf(section_data, body_style))
                story.append(Spacer(1, 0.15 * inch))

            doc.build(story)
            with open("/tmp/decisionlens_report.pdf", "rb") as f:
                return f.read()
        except Exception as e:
            logger.error("[PDFExporter] %s", e)
            raise RuntimeError(f"PDF export failed: {e}")

    @staticmethod
    def _render_section_pdf(data: Any, style) -> List[Any]:
        from reportlab.platypus import Paragraph, Spacer, Table, TableStyle
        from reportlab.lib import colors
        from reportlab.lib.styles import ParagraphStyle

        elements = []
        if data is None:
            return elements
        if isinstance(data, str):
            elements.append(Paragraph(data.replace("\n", "<br/>"), style))
            return elements
        if isinstance(data, (int, float)):
            elements.append(Paragraph(str(data), style))
            return elements
        if isinstance(data, list):
            for item in data[:20]:
                if isinstance(item, dict):
                    elements.extend(PDFExporter._render_dict_pdf(item, style))
                else:
                    elements.append(Paragraph(f"- {str(item)}", style))
            return elements
        if isinstance(data, dict):
            elements.extend(PDFExporter._render_dict_pdf(data, style))
            return elements
        elements.append(Paragraph(str(data), style))
        return elements

    @staticmethod
    def _render_dict_pdf(d: Dict[str, Any], style) -> List[Any]:
        from reportlab.platypus import Paragraph, Spacer, Table, TableStyle
        from reportlab.lib import colors
        from reportlab.lib.styles import ParagraphStyle

        elements = []
        bullet_style = ParagraphStyle(
            "Bullet",
            parent=style,
            leftIndent=12,
            bulletIndent=6,
            spaceAfter=4,
        )
        for key, value in d.items():
            if value is None:
                continue
            if isinstance(value, (dict, list)):
                header = f"<b>{key}:</b>"
                elements.append(Paragraph(header, style))
                elements.extend(PDFExporter._render_section_pdf(value, style))
            else:
                line = f"<b>{key}:</b> {value}"
                elements.append(Paragraph(line, bullet_style))
        return elements


class DOCXExporter:
    @staticmethod
    def export(report: Dict[str, Any], title: str = "DecisionLens Executive Report") -> bytes:
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
            from openpyxl.utils import get_column_letter

            wb = Workbook()
            ws = wb.active
            ws.title = "Executive Report"

            header_font = Font(name="Calibri", size=16, bold=True, color="FFFFFF")
            subheader_font = Font(name="Calibri", size=12, bold=True, color="1a1a2e")
            body_font = Font(name="Calibri", size=10)
            header_fill = PatternFill(start_color="1a1a2e", end_color="1a1a2e", fill_type="solid")
            section_fill = PatternFill(start_color="e94560", end_color="e94560", fill_type="solid")
            thin_border = Border(
                left=Side(style="thin", color="CCCCCC"),
                right=Side(style="thin", color="CCCCCC"),
                top=Side(style="thin", color="CCCCCC"),
                bottom=Side(style="thin", color="CCCCCC"),
            )

            ws.merge_cells("A1:F1")
            ws["A1"] = title
            ws["A1"].font = header_font
            ws["A1"].fill = header_fill
            ws["A1"].alignment = Alignment(horizontal="center", vertical="center")
            ws.row_dimensions[1].height = 36

            meta_row = 3
            ws.cell(row=meta_row, column=1, value="Generated:").font = body_font
            ws.cell(row=meta_row, column=2, value=report.get("generated_at", "")[:19])
            ws.cell(row=meta_row + 1, column=1, value="Domain:").font = body_font
            ws.cell(row=meta_row + 1, column=2, value=report.get("domain", "N/A"))
            ws.cell(row=meta_row + 2, column=1, value="Dataset Type:").font = body_font
            ws.cell(row=meta_row + 2, column=2, value=report.get("dataset_type", "N/A"))

            sections = report.get("sections", {})
            section_order = [
                "executive_summary", "business_health", "kpi_summary", "kpi_overview",
                "revenue", "customers", "products", "categories", "trend_analysis",
                "forecast", "root_cause_analysis", "risks", "opportunities",
                "recommendations", "evidence", "charts", "what_if_analysis",
                "recommended_actions", "business_impact", "key_findings",
                "confidence_evidence", "roadmap_30_90_180", "domain_specific",
            ]

            current_row = meta_row + 4
            for section_name in section_order:
                section_data = sections.get(section_name)
                if section_data is None:
                    continue

                ws.merge_cells(start_row=current_row, start_column=1, end_row=current_row, end_column=6)
                cell = ws.cell(row=current_row, column=1, value=section_name.replace("_", " ").title())
                cell.font = subheader_font
                cell.fill = section_fill
                cell.alignment = Alignment(horizontal="left", vertical="center")
                cell.font = Font(name="Calibri", size=12, bold=True, color="FFFFFF")
                ws.row_dimensions[current_row].height = 24
                current_row += 1

                rows_written = DOCXExporter._write_section_xlsx(ws, current_row, section_data, body_font, thin_border)
                current_row += rows_written + 1

            for col in range(1, 7):
                ws.column_dimensions[get_column_letter(col)].width = 22

            output_path = "/tmp/decisionlens_report.xlsx"
            wb.save(output_path)
            with open(output_path, "rb") as f:
                return f.read()
        except Exception as e:
            logger.error("[DOCXExporter] %s", e)
            raise RuntimeError(f"DOCX export failed: {e}")

    @staticmethod
    def _write_section_xlsx(ws, start_row: int, data: Any, font, border) -> int:
        rows = 0
        if data is None:
            return 0
        if isinstance(data, str):
            ws.cell(row=start_row, column=1, value=data).font = font
            return 1
        if isinstance(data, (int, float)):
            ws.cell(row=start_row, column=1, value=str(data)).font = font
            return 1
        if isinstance(data, list):
            for item in data[:15]:
                if isinstance(item, dict):
                    rows += DOCXExporter._write_dict_xlsx(ws, start_row + rows, item, font, border)
                else:
                    ws.cell(row=start_row + rows, column=1, value=str(item)).font = font
                    rows += 1
            return rows
        if isinstance(data, dict):
            return DOCXExporter._write_dict_xlsx(ws, start_row, data, font, border)
        ws.cell(row=start_row, column=1, value=str(data)).font = font
        return 1

    @staticmethod
    def _write_dict_xlsx(ws, start_row: int, d: Dict[str, Any], font, border) -> int:
        rows = 0
        for key, value in d.items():
            if value is None:
                continue
            row = start_row + rows
            ws.cell(row=row, column=1, value=str(key)).font = Font(bold=True)
            ws.cell(row=row, column=1).border = border
            if isinstance(value, (dict, list)):
                nested_rows = DOCXExporter._write_section_xlsx(ws, row, value, font, border)
                rows = max(rows, nested_rows)
            else:
                ws.cell(row=row, column=2, value=str(value)).font = font
                ws.cell(row=row, column=2).border = border
                rows += 1
        return rows


class PPTXExporter:
    @staticmethod
    def export(report: Dict[str, Any], title: str = "DecisionLens Executive Report") -> bytes:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RgbColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.enum.shapes import MSO_SHAPE

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            blank_layout = prs.slide_layouts[6]

            slide = prs.slides.add_slide(blank_layout)
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
            tf = title_box.text_frame
            tf.text = title
            p = tf.paragraphs[0]
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = RgbColor(0x1A, 0x1A, 0x2E)
            p.alignment = PP_ALIGN.LEFT

            meta_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(1))
            meta_tf = meta_box.text_frame
            meta_tf.text = (
                f"Generated: {report.get('generated_at', '')[:19]}  |  "
                f"Domain: {report.get('domain', 'N/A')}  |  "
                f"Dataset: {report.get('dataset_type', 'N/A')}"
            )
            for paragraph in meta_tf.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = RgbColor(0x66, 0x66, 0x66)

            sections = report.get("sections", {})
            section_order = [
                "executive_summary", "business_health", "kpi_summary", "kpi_overview",
                "revenue", "customers", "products", "categories", "trend_analysis",
                "forecast", "root_cause_analysis", "risks", "opportunities",
                "recommendations", "evidence", "charts", "what_if_analysis",
                "recommended_actions", "business_impact", "key_findings",
                "confidence_evidence", "roadmap_30_90_180", "domain_specific",
            ]

            for section_name in section_order:
                section_data = sections.get(section_name)
                if section_data is None:
                    continue

                slide = prs.slides.add_slide(blank_layout)
                header = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
                htf = header.text_frame
                htf.text = section_name.replace("_", " ").title()
                hp = htf.paragraphs[0]
                hp.font.size = Pt(24)
                hp.font.bold = True
                hp.font.color.rgb = RgbColor(0xE9, 0x45, 0x60)

                content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12), Inches(6))
                ctf = content_box.text_frame
                ctf.word_wrap = True

                lines = PPTXExporter._flatten_section(section_data)
                for i, line in enumerate(lines[:20]):
                    if i == 0:
                        p = ctf.paragraphs[0]
                    else:
                        p = ctf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(11)
                    p.font.color.rgb = RgbColor(0x33, 0x33, 0x33)
                    p.space_after = Pt(6)

            output_path = "/tmp/decisionlens_report.pptx"
            prs.save(output_path)
            with open(output_path, "rb") as f:
                return f.read()
        except Exception as e:
            logger.error("[PPTXExporter] %s", e)
            raise RuntimeError(f"PPTX export failed: {e}")

    @staticmethod
    def _flatten_section(data: Any, max_depth: int = 3) -> List[str]:
        lines = []
        if max_depth <= 0:
            return lines
        if data is None:
            return lines
        if isinstance(data, str):
            return [data]
        if isinstance(data, (int, float)):
            return [str(data)]
        if isinstance(data, list):
            for item in data[:10]:
                lines.extend(PPTXExporter._flatten_section(item, max_depth - 1))
            return lines
        if isinstance(data, dict):
            for key, value in data.items():
                if value is None:
                    continue
                if isinstance(value, (dict, list)):
                    lines.append(f"{key}:")
                    lines.extend(PPTXExporter._flatten_section(value, max_depth - 1))
                else:
                    lines.append(f"{key}: {value}")
        return lines
